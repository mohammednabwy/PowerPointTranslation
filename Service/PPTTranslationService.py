import os
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.enum.lang import MSO_LANGUAGE_ID
from Service import Translator
#import pptx
#from docx.oxml.shared import qn
#---------------------------------------------
def savePowerPoint(filename,presentation):
    presentation.save(filename)
    return
#---------------------------------------------
def translatePPTFile(pptfilePath,target='ar'):     
    presentation = Presentation(pptfilePath)
    for slideIndex in range(0,len(presentation.slides)): 
       print("Translating slide number : ", str(slideIndex+1))
       slide=presentation.slides[slideIndex]  
     
       for i in range(0,len(slide.shapes)):        
           shape=slide.shapes[i] 
           if shape.has_text_frame:             
               text_frame = shape.text_frame
               for paragraph in text_frame.paragraphs:                  
#                    if target=='ar':
#                        pPr=paragraph._p.get_or_add_pPr()                       
#                        ap= pptx.oxml.xmlchemy.OxmlElement('a:p')
#                        ap.set(qn('a:pPr'),'rtl')
#                        pPr.append(ap)
                    for run in paragraph.runs:
                        try:
                            cur_text = run.text
                            if len(cur_text) < 3 : continue
                            new_text = Translator.translate(cur_text,target) 
                            run.text = new_text 
                            if target=='ar':
                                run.alignment = PP_ALIGN.RIGHT
                                run.font.language_id =MSO_LANGUAGE_ID.ARABIC
                                
                        except:
                            continue
           if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells: 
                            if len(cell.text) < 3 : continue
                            new_text = Translator.translate(cell.text,target)
                            cell.text = new_text
    #save 
    file_name=os.path.basename(pptfilePath)
    file_directory=os.path.dirname(pptfilePath) 
    translated_file_name='translated_' + target +"_" +  file_name.split('.')[0] + '.'  +file_name.split('.')[1]
    translated_file = os.path.join(file_directory, translated_file_name)
    savePowerPoint(translated_file,presentation)
    return translated_file
#---------------------------------------------
def main():
    pptfilePath='Files\\DataStructureAndAlgorthimDesign_1.pptx'  
    translatePPTFile(pptfilePath,target='ar')   
    return